"""PowerPoint (PPTX) generation, editing, and reading tools.

Uses python-pptx to create, edit, and extract content from
PowerPoint presentations programmatically.

License: MIT
Version: 1.0
Category: productivity
"""

from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from tools import tool

# ──────────────────────────────────────────────
# Color palettes / themes
# ──────────────────────────────────────────────

THEMES = {
    "professional": {
        "primary": RGBColor(0x1A, 0x3C, 0x6E),      # Dark blue
        "secondary": RGBColor(0x2E, 0x86, 0xDE),     # Bright blue
        "accent": RGBColor(0xE6, 0x7E, 0x22),        # Orange
        "background": RGBColor(0xF8, 0xF9, 0xFA),    # Light gray
        "text": RGBColor(0x2C, 0x3E, 0x50),          # Dark gray
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xEB, 0xF5, 0xFB),     # Light blue bg
    },
    "modern": {
        "primary": RGBColor(0x2C, 0x3E, 0x50),
        "secondary": RGBColor(0x18, 0xBC, 0x9C),     # Teal
        "accent": RGBColor(0xE7, 0x4C, 0x3C),        # Red
        "background": RGBColor(0xEC, 0xF0, 0xF1),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xE8, 0xF8, 0xF5),
    },
    "colorful": {
        "primary": RGBColor(0x8E, 0x44, 0xAD),       # Purple
        "secondary": RGBColor(0x34, 0x98, 0xDB),     # Blue
        "accent": RGBColor(0xE6, 0x7E, 0x22),        # Orange
        "background": RGBColor(0xF5, 0xF5, 0xF5),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xF3, 0xE5, 0xF5),
    },
    "minimal": {
        "primary": RGBColor(0x33, 0x33, 0x33),
        "secondary": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0x00, 0x7A, 0xFF),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x33, 0x33, 0x33),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xF5, 0xF5, 0xF5),
    },
}

DEFAULT_THEME = "professional"
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)


def _get_theme(name: str | None = None) -> dict:
    """Get a theme dict by name, falling back to default."""
    return THEMES.get(name or DEFAULT_THEME, THEMES[DEFAULT_THEME])


def _new_slide(prs):
    """Create a new blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=None, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_separator(slide, left, top, width, color):
    """Add a thin colored separator line."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_bullet_slide(prs, title_text, bullets, theme):
    """Add a content slide with bullet points."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  title_text, font_size=32, bold=True, color=theme["primary"])
    _add_separator(slide, Inches(0.8), Inches(1.3), Inches(2), theme["secondary"])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.font.color.rgb = theme["text"]
        p.space_after = Pt(12)


def _add_cover_slide(prs, title, subtitle, theme):
    """Add a cover/title slide."""
    slide = _new_slide(prs)

    # Background color bar (top half)
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(4.5),
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = theme["primary"]
    bg_bar.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SLIDE_WIDTH, Inches(0.08),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = theme["accent"]
    stripe.line.fill.background()

    # Title
    _add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5),
                  title, font_size=44, bold=True, color=theme["white"],
                  alignment=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        _add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                      subtitle, font_size=24, color=RGBColor(0xDD, 0xDD, 0xDD),
                      alignment=PP_ALIGN.LEFT)

    # Bottom info
    _add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
                  "Generated with pptx-generator v1.0", font_size=14,
                  color=theme["text"], alignment=PP_ALIGN.LEFT)


def _add_toc_slide(prs, items, theme):
    """Add a table of contents slide."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  "📋 目录 / Table of Contents", font_size=32, bold=True,
                  color=theme["primary"])
    _add_separator(slide, Inches(0.8), Inches(1.3), Inches(2), theme["secondary"])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}.  {item}"
        p.font.size = Pt(22)
        p.font.name = "Calibri"
        p.font.color.rgb = theme["text"]
        p.space_after = Pt(14)


def _add_section_slide(prs, section_title, theme):
    """Add a section divider slide."""
    slide = _new_slide(prs)

    # Full background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary"]
    bg.line.fill.background()

    # Accent line
    _add_separator(slide, Inches(1.5), Inches(3.2), Inches(3), theme["accent"])

    # Section title
    _add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
                  section_title, font_size=40, bold=True, color=theme["white"],
                  alignment=PP_ALIGN.LEFT)


def _add_summary_slide(prs, title, points, theme):
    """Add a summary/conclusion slide."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  title or "📌 总结 / Summary", font_size=32, bold=True,
                  color=theme["primary"])
    _add_separator(slide, Inches(0.8), Inches(1.3), Inches(2), theme["accent"])

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✅  {point}"
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.font.color.rgb = theme["text"]
        p.space_after = Pt(12)


def _add_two_column_slide(prs, title, left_content, right_content, theme):
    """Add a two-column layout slide."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  title, font_size=32, bold=True, color=theme["primary"])
    _add_separator(slide, Inches(0.8), Inches(1.3), Inches(2), theme["secondary"])

    # Left column
    txBox_left = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
    )
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    for i, line_text in enumerate(left_content):
        p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        p.text = line_text
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.font.color.rgb = theme["text"]
        p.space_after = Pt(8)

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(0.04), Inches(4.5),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = theme["light_bg"]
    div.line.fill.background()

    # Right column
    txBox_right = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0),
    )
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    for i, line_text in enumerate(right_content):
        p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
        p.text = line_text
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.font.color.rgb = theme["text"]
        p.space_after = Pt(8)


def _add_table_slide(prs, title, headers, rows, theme):
    """Add a slide with a table."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  title, font_size=32, bold=True, color=theme["primary"])

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_height = max(Inches(0.4 * num_rows + 0.3), Inches(2.0))

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.5), Inches(11.7), table_height,
    )
    table = table_shape.table

    # Style header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["primary"]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme["white"]
            paragraph.font.name = "Calibri"

    # Style data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme["light_bg"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = theme["text"]
                paragraph.font.name = "Calibri"


def _add_content_slide(prs, title, body_text, theme):
    """Add a standard content slide with body text."""
    slide = _new_slide(prs)
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
                  title, font_size=32, bold=True, color=theme["primary"])
    _add_separator(slide, Inches(0.8), Inches(1.3), Inches(2), theme["secondary"])
    _add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5),
                  body_text, font_size=20, color=theme["text"])


# ──────────────────────────────────────────────
# Public tools
# ──────────────────────────────────────────────


@tool(
    name="pptx_list_templates",
    description="List available PowerPoint presentation templates/themes.",
    parameters={
        "type": "object",
        "properties": {},
        "required": [],
    },
)
def pptx_list_templates() -> str:
    """List available presentation themes."""
    themes_info = []
    for name, colors in THEMES.items():
        themes_info.append(
            f"• {name}: primary=#{colors['primary']}, "
            f"secondary=#{colors['secondary']}, accent=#{colors['accent']}"
        )
    return (
        "Available presentation themes:\n" + "\n".join(themes_info) + "\n\n"
        "Slide types supported: cover, toc, content, section, summary, "
        "two_column, table"
    )


@tool(
    name="pptx_create",
    description="Create a new PowerPoint presentation from scratch. "
                "Supports cover, TOC, content, section divider, summary, "
                "two_column, and table slides.",
    parameters={
        "type": "object",
        "properties": {
            "output_path": {
                "type": "string",
                "description": "Absolute path to save the PPTX file (e.g., /path/to/presentation.pptx)",
            },
            "title": {
                "type": "string",
                "description": "Presentation title (used on cover slide)",
            },
            "subtitle": {
                "type": "string",
                "description": "Optional subtitle for the cover slide",
            },
            "slides": {
                "type": "string",
                "description": "JSON array of slide definitions. Each slide has: "
                               "{type, title, content, items, bullets, left_content, "
                               "right_content, headers, rows, points}. "
                               "Types: cover, toc, content, section, summary, two_column, table. "
                               "Example: [{\"type\":\"cover\"}, {\"type\":\"toc\",\"items\":[\"Intro\",\"Body\"]}, "
                               "{\"type\":\"content\",\"title\":\"Slide 1\",\"content\":\"Hello\"}]",
            },
            "theme": {
                "type": "string",
                "description": "Theme name: professional, modern, colorful, or minimal (default: professional)",
                "enum": ["professional", "modern", "colorful", "minimal"],
            },
        },
        "required": ["output_path", "title", "slides"],
    },
)
def pptx_create(output_path: str, title: str, subtitle: str = "",
                slides: str = "[]", theme: str = "professional") -> str:
    """Create a PowerPoint presentation from slide definitions."""
    try:
        theme_colors = _get_theme(theme)
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        slide_list = json.loads(slides) if isinstance(slides, str) else slides

        for i, slide_def in enumerate(slide_list):
            stype = slide_def.get("type", "content")

            if stype == "cover":
                _add_cover_slide(
                    prs,
                    slide_def.get("title", title),
                    slide_def.get("subtitle", subtitle),
                    theme_colors,
                )
            elif stype == "toc":
                _add_toc_slide(
                    prs,
                    slide_def.get("items", slide_def.get("bullets", [])),
                    theme_colors,
                )
            elif stype == "section":
                _add_section_slide(
                    prs,
                    slide_def.get("title", "Section"),
                    theme_colors,
                )
            elif stype == "summary":
                _add_summary_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("points", slide_def.get("items", [])),
                    theme_colors,
                )
            elif stype == "two_column":
                _add_two_column_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("left_content", []),
                    slide_def.get("right_content", []),
                    theme_colors,
                )
            elif stype == "table":
                _add_table_slide(
                    prs,
                    slide_def.get("title", "Table"),
                    slide_def.get("headers", []),
                    slide_def.get("rows", []),
                    theme_colors,
                )
            else:  # content (default)
                content_text = slide_def.get("content", "")
                bullets = slide_def.get("bullets", slide_def.get("items", []))
                if bullets:
                    _add_bullet_slide(
                        prs,
                        slide_def.get("title", f"Slide {i+1}"),
                        bullets,
                        theme_colors,
                    )
                else:
                    _add_content_slide(
                        prs,
                        slide_def.get("title", f"Slide {i+1}"),
                        content_text,
                        theme_colors,
                    )

        # Save
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output))

        slide_count = len(prs.slides)
        return json.dumps({
            "success": True,
            "path": str(output.resolve()),
            "slides": slide_count,
            "file_size": f"{output.stat().st_size / 1024:.1f} KB",
            "message": f"Presentation created successfully: {slide_count} slides, "
                       f"theme: {theme}",
        }, ensure_ascii=False)

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to create presentation: {type(e).__name__}: {e}",
        }, ensure_ascii=False)


@tool(
    name="pptx_edit",
    description="Edit an existing PPTX file. Can modify text on slides, "
                "add new slides, or change the presentation metadata.",
    parameters={
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the existing PPTX file",
            },
            "action": {
                "type": "string",
                "description": "Action to perform: 'add_slides' (append new slides), "
                               "'replace_text' (find and replace text in all slides), "
                               "'set_title' (change slide title by index)",
                "enum": ["add_slides", "replace_text", "set_title"],
            },
            "slides": {
                "type": "string",
                "description": "JSON array of slide definitions (for add_slides action). "
                               "Same format as pptx_create slides parameter.",
            },
            "old_text": {
                "type": "string",
                "description": "Text to find (for replace_text action)",
            },
            "new_text": {
                "type": "string",
                "description": "Replacement text (for replace_text action)",
            },
            "slide_index": {
                "type": "integer",
                "description": "0-based slide index (for set_title action)",
            },
            "new_title": {
                "type": "string",
                "description": "New title text (for set_title action)",
            },
            "theme": {
                "type": "string",
                "description": "Theme for new slides: professional, modern, colorful, minimal",
                "enum": ["professional", "modern", "colorful", "minimal"],
            },
        },
        "required": ["file_path", "action"],
    },
)
def pptx_edit(file_path: str, action: str = "add_slides",
              slides: str = "[]", old_text: str = "", new_text: str = "",
              slide_index: int = 0, new_title: str = "",
              theme: str = "professional") -> str:
    """Edit an existing PowerPoint presentation."""
    try:
        path = Path(file_path)
        if not path.exists():
            return json.dumps({
                "success": False,
                "error": f"File not found: {file_path}",
            })

        prs = Presentation(str(path))
        theme_colors = _get_theme(theme)

        if action == "add_slides":
            slide_list = json.loads(slides) if isinstance(slides, str) else slides
            for slide_def in slide_list:
                stype = slide_def.get("type", "content")
                if stype == "cover":
                    _add_cover_slide(prs, slide_def.get("title", ""),
                                     slide_def.get("subtitle", ""), theme_colors)
                elif stype == "toc":
                    _add_toc_slide(prs, slide_def.get("items", []), theme_colors)
                elif stype == "section":
                    _add_section_slide(prs, slide_def.get("title", ""), theme_colors)
                elif stype == "summary":
                    _add_summary_slide(prs, slide_def.get("title", ""),
                                       slide_def.get("points", []), theme_colors)
                elif stype == "two_column":
                    _add_two_column_slide(prs, slide_def.get("title", ""),
                                          slide_def.get("left_content", []),
                                          slide_def.get("right_content", []), theme_colors)
                elif stype == "table":
                    _add_table_slide(prs, slide_def.get("title", ""),
                                     slide_def.get("headers", []),
                                     slide_def.get("rows", []), theme_colors)
                else:
                    _add_content_slide(prs, slide_def.get("title", ""),
                                       slide_def.get("content", ""), theme_colors)

            prs.save(str(path))
            return json.dumps({
                "success": True,
                "message": f"Added {len(slide_list)} slide(s) to {file_path}",
                "total_slides": len(prs.slides),
            })

        elif action == "replace_text":
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    count += 1
            prs.save(str(path))
            return json.dumps({
                "success": True,
                "message": f"Replaced '{old_text}' with '{new_text}' in {count} text runs",
            })

        elif action == "set_title":
            if slide_index < 0 or slide_index >= len(prs.slides):
                return json.dumps({
                    "success": False,
                    "error": f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})",
                })
            slide = prs.slides[slide_index]
            if slide.shapes.title:
                slide.shapes.title.text = new_title
            else:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape.text_frame.paragraphs[0].text = new_title
                        break
            prs.save(str(path))
            return json.dumps({
                "success": True,
                "message": f"Set title of slide {slide_index} to '{new_title}'",
            })

        else:
            return json.dumps({
                "success": False,
                "error": f"Unknown action: {action}",
            })

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to edit presentation: {type(e).__name__}: {e}",
        }, ensure_ascii=False)


@tool(
    name="pptx_read",
    description="Extract text content from an existing PPTX file. "
                "Returns slide titles and body text for each slide.",
    parameters={
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Absolute path to the PPTX file to read",
            },
        },
        "required": ["file_path"],
    },
)
def pptx_read(file_path: str) -> str:
    """Read and extract text content from a PowerPoint file."""
    try:
        path = Path(file_path)
        if not path.exists():
            return json.dumps({
                "success": False,
                "error": f"File not found: {file_path}",
            })

        prs = Presentation(str(path))
        result = {
            "success": True,
            "file": str(path.resolve()),
            "slide_count": len(prs.slides),
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
            "slides": [],
        }

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "index": i,
                "texts": [],
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_info["texts"].append(text)
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_info["table"] = table_data
            result["slides"].append(slide_info)

        return json.dumps(result, ensure_ascii=False, indent=2)

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to read presentation: {type(e).__name__}: {e}",
        }, ensure_ascii=False)
